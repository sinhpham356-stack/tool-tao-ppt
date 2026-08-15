import io
import re
import math
import os
import urllib.request
from datetime import datetime
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps, ImageDraw, ImageFont

# --- SỬA LỖI ĐỌC ẢNH IPHONE (HEIC) ---
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass

st.set_page_config(page_title="Công Cụ Chèn Ảnh PowerPoint & PDF", page_icon="⚡", layout="centered")

# ==========================================
# KHỞI TẠO BỘ NHỚ TẠM
# ==========================================
if "photo_cart" not in st.session_state: st.session_state.photo_cart = {} 
if "template_bytes" not in st.session_state: st.session_state.template_bytes = None 
if "cover_pos" not in st.session_state: st.session_state.cover_pos = "Dưới - Giữa"
if "ty_pos" not in st.session_state: st.session_state.ty_pos = "Trung tâm"

# Găm file vào bộ nhớ để không bị mất khi bấm Tải xuống
if "final_pdf" not in st.session_state: st.session_state.final_pdf = None
if "final_pptx" not in st.session_state: st.session_state.final_pptx = None
if "show_download_pdf" not in st.session_state: st.session_state.show_download_pdf = False
if "show_download_pptx" not in st.session_state: st.session_state.show_download_pptx = False

# ==========================================
# HÀM LÕI XỬ LÝ CHUNG
# ==========================================
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def process_bg_image(uploaded_file, target_ratio=16/9):
    if not uploaded_file: return None
    try:
        bg_bytes = uploaded_file.getvalue()
        with Image.open(io.BytesIO(bg_bytes)) as img:
            img = ImageOps.exif_transpose(img)
            if img.mode != 'RGB': img = img.convert('RGB')
            img_ratio = img.width / img.height
            if img_ratio > target_ratio:
                new_w = int(img.height * target_ratio)
                offset = (img.width - new_w) // 2
                img = img.crop((offset, 0, offset + new_w, img.height))
            elif img_ratio < target_ratio:
                new_h = int(img.width / target_ratio)
                offset = (img.height - new_h) // 2
                img = img.crop((0, offset, img.width, offset + new_h))
            bg_stream = io.BytesIO()
            img.save(bg_stream, format='JPEG', quality=95)
            return bg_stream
    except Exception: return None

# ------------------------------------------
# HÀM VẼ PPTX
# ------------------------------------------
def add_image_exact(slide, img_stream, left, top, width, height):
    with Image.open(img_stream) as img:
        img = ImageOps.exif_transpose(img) 
        if img.mode != 'RGB': img = img.convert('RGB')
        temp_stream = io.BytesIO()
        img.save(temp_stream, format='JPEG', quality=95)
        temp_stream.seek(0)
    pic = slide.shapes.add_picture(temp_stream, int(left), int(top), width=int(width), height=int(height))
    pic.line.color.rgb = RGBColor(30, 30, 30)
    pic.line.width = Inches(0.03) 

def move_slide(prs, old_index, new_index):
    if old_index == new_index: return
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_to_move = slides[old_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(new_index, slide_to_move)

def draw_adaptive_grid(slide, layout_rows, start_x_base, start_y_base, usable_w, usable_h, GAP):
    num_rows = len(layout_rows)
    if num_rows == 0: return
    H_max = (usable_h - (num_rows - 1) * GAP) / num_rows
    H_final = H_max
    for row in layout_rows:
        if not row: continue
        sum_ratios = sum(img['w'] / img['h'] for img in row)
        H_row_max = (usable_w - (len(row) - 1) * GAP) / sum_ratios
        if H_row_max < H_final: H_final = H_row_max
    Total_H = num_rows * H_final + (num_rows - 1) * GAP
    current_y = start_y_base + (usable_h - Total_H) / 2
    for row in layout_rows:
        if not row: continue
        row_w = sum(H_final * (img['w'] / img['h']) for img in row) + (len(row) - 1) * GAP
        current_x = start_x_base + (usable_w - row_w) / 2
        for img in row:
            img_w = H_final * (img['w'] / img['h'])
            add_image_exact(slide, img['stream'], current_x, current_y, img_w, H_final)
            current_x += img_w + GAP
        current_y += H_final + GAP

# => ĐÂY LÀ HÀM BỊ LỠ TAY XÓA MẤT (ĐÃ THÊM LẠI)
def partition_images(imgs, max_size):
    if not imgs: return []
    n = len(imgs)
    num_slides = math.ceil(n / max_size)
    base = n // num_slides
    rem = n % num_slides
    res = []
    idx = 0
    for i in range(num_slides):
        s = base + 1 if i < rem else base
        res.append(imgs[idx:idx+s])
        idx += s
    return res

# ------------------------------------------
# HÀM VẼ PDF ĐỘC LẬP BẰNG PILLOW
# ------------------------------------------
def emu_to_px(emu):
    return int((emu / 914400.0) * 120)

def add_pdf_slide(pdf_slides, bg_stream=None):
    img = Image.new('RGB', (1600, 900), (255, 255, 255))
    if bg_stream:
        bg_stream.seek(0)
        bg = Image.open(bg_stream).convert('RGB')
        bg = bg.resize((1600, 900), Image.LANCZOS)
        img.paste(bg, (0, 0))
    pdf_slides.append(img)
    return img

def draw_text_pdf(slide_img, text, x_emu, y_emu, w_emu, h_emu, align_pos, pt_size, color_hex, underline=False):
    draw = ImageDraw.Draw(slide_img)
    px, py, pw = emu_to_px(x_emu), emu_to_px(y_emu), emu_to_px(w_emu)
    font_size = int(pt_size * 120 / 72)
    
    font_path = "Roboto-Medium.ttf"
    if not os.path.exists(font_path):
        try:
            urllib.request.urlretrieve("https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Medium.ttf", font_path)
        except: pass
    try: font = ImageFont.truetype(font_path, font_size)
    except: font = ImageFont.load_default()
        
    r, g, b = hex_to_rgb(color_hex)
    
    words = text.split()
    lines = []
    current_line = ""
    for word in words:
        test_line = current_line + " " + word if current_line else word
        bbox = draw.textbbox((0,0), test_line, font=font)
        if bbox[2] - bbox[0] <= pw:
            current_line = test_line
        else:
            if current_line: lines.append(current_line)
            current_line = word
    if current_line: lines.append(current_line)
        
    current_y = py
    for line in lines:
        bbox = draw.textbbox((0,0), line, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        
        if "Trái" in align_pos or align_pos == PP_ALIGN.LEFT: tx = px
        elif "Phải" in align_pos or align_pos == PP_ALIGN.RIGHT: tx = px + pw - tw
        else: tx = px + (pw - tw) / 2
            
        draw.text((tx, current_y), line, font=font, fill=(r,g,b))
        if underline:
            draw.line([tx, current_y + th + 5, tx + tw, current_y + th + 5], fill=(r,g,b), width=3)
        current_y += th + 15

def add_image_pdf(slide_img, img_stream, left_emu, top_emu, width_emu, height_emu):
    img_stream.seek(0)
    img = Image.open(img_stream).convert('RGB')
    px, py = emu_to_px(left_emu), emu_to_px(top_emu)
    pw, ph = emu_to_px(width_emu), emu_to_px(height_emu)
    img = img.resize((pw, ph), Image.LANCZOS)
    slide_img.paste(img, (px, py))
    draw = ImageDraw.Draw(slide_img)
    draw.rectangle([px, py, px+pw, py+ph], outline=(30,30,30), width=4)

def draw_adaptive_grid_pdf(slide_img, layout_rows, start_x_base, start_y_base, usable_w, usable_h, GAP):
    num_rows = len(layout_rows)
    if num_rows == 0: return
    H_max = (usable_h - (num_rows - 1) * GAP) / num_rows
    H_final = H_max
    for row in layout_rows:
        if not row: continue
        sum_ratios = sum(img['w'] / img['h'] for img in row)
        H_row_max = (usable_w - (len(row) - 1) * GAP) / sum_ratios
        if H_row_max < H_final: H_final = H_row_max
    Total_H = num_rows * H_final + (num_rows - 1) * GAP
    current_y = start_y_base + (usable_h - Total_H) / 2
    for row in layout_rows:
        if not row: continue
        row_w = sum(H_final * (img['w'] / img['h']) for img in row) + (len(row) - 1) * GAP
        current_x = start_x_base + (usable_w - row_w) / 2
        for img in row:
            img_w = H_final * (img['w'] / img['h'])
            add_image_pdf(slide_img, img['stream'], current_x, current_y, img_w, H_final)
            current_x += img_w + GAP
        current_y += H_final + GAP

# ----------------- 9 ĐIỂM NEO VỊ TRÍ -----------------
def render_position_grid(state_key, key_prefix):
    def set_pos(pos): st.session_state[state_key] = pos
    current_pos = st.session_state[state_key]
    grid_config = [
        [("Trên - Trái", "↖️"), ("Trên - Giữa", "⬆️"), ("Trên - Phải", "↗️")],
        [("Giữa - Trái", "⬅️"), ("Trung tâm", "⏺️"), ("Giữa - Phải", "➡️")],
        [("Dưới - Trái", "↙️"), ("Dưới - Giữa", "⬇️"), ("Dưới - Phải", "↘️")]
    ]
    st.markdown(f"📍 Đang định vị: **{current_pos}**")
    for r_idx, row in enumerate(grid_config):
        cols = st.columns(3, gap="small")
        for c_idx, (pos_name, icon) in enumerate(row):
            btn_label = f"✅ {icon}" if current_pos == pos_name else icon
            cols[c_idx].button(
                btn_label, key=f"{key_prefix}_{r_idx}_{c_idx}", 
                use_container_width=True, on_click=set_pos, args=(pos_name,)
            )

def get_alignment(pos_str):
    if "Trái" in pos_str: return PP_ALIGN.LEFT
    elif "Phải" in pos_str: return PP_ALIGN.RIGHT
    else: return PP_ALIGN.CENTER

def get_cover_y(pos_str, slide_h, is_sub=False):
    if "Trên" in pos_str: return Inches(0.5) if not is_sub else Inches(1.5)
    elif "Dưới" in pos_str: return slide_h - Inches(2.2) if not is_sub else slide_h - Inches(1.2)
    else: return slide_h/2 - Inches(1.0) if not is_sub else slide_h/2 + Inches(0.2)

def get_ty_y(pos_str, slide_h):
    if "Trên" in pos_str: return Inches(0.8)
    elif "Dưới" in pos_str: return slide_h - Inches(1.8)
    else: return slide_h/2 - Inches(0.5)

# ==========================================
# GIAO DIỆN HIỂN THỊ
# ==========================================
st.title("⚡ TRỢ LÝ TẠO REPORT BẰNG HÌNH ẢNH")
st.error("🍎 **CHỐNG VĂNG ẢNH IPHONE:** Tải lên từng đợt nhỏ (2-5 ảnh), hệ thống tự gom vào Giỏ Hàng!")

st.header("Bước 1: Tải lên File PowerPoint Mẫu (Không bắt buộc)")
template_file = st.file_uploader("Chọn file .pptx (Chỉ cần up 1 lần)", type=["pptx"])
if template_file: st.session_state.template_bytes = template_file.getvalue()
if st.session_state.template_bytes:
    st.success("✅ Đã lưu file PPTX mẫu vào bộ nhớ!")
    if st.button("🗑️ Xóa file mẫu (Để tạo file trắng mới)"):
        st.session_state.template_bytes = None
        st.rerun()

st.header("Bước 2: Ném ảnh vào Giỏ")
uploaded_images = st.file_uploader("📂 Nhấn để chọn ảnh từ máy", accept_multiple_files=True)
if uploaded_images:
    count = 0
    for f in uploaded_images:
        if f.name not in st.session_state.photo_cart:
            st.session_state.photo_cart[f.name] = {"bytes": f.getvalue(), "name": f.name}
            count += 1
    if count > 0: st.success(f"🎉 Vừa nhặt thêm {count} ảnh vào giỏ!")

enable_camera = st.toggle("📷 Bật máy ảnh để chụp trực tiếp")
if enable_camera:
    camera_photo = st.camera_input("Chụp ảnh thực tế tại hiện trường")
    if camera_photo:
        cam_bytes = camera_photo.getvalue()
        is_duplicate = False
        for item in st.session_state.photo_cart.values():
            if item["bytes"] == cam_bytes:
                is_duplicate = True; break
        if not is_duplicate:
            cam_name = f"cam_{datetime.now().strftime('%H%M%S')}.jpg"
            st.session_state.photo_cart[cam_name] = {"bytes": cam_bytes, "name": cam_name}
            st.success("📸 Đã ném ảnh vừa chụp vào giỏ!")

if st.session_state.photo_cart:
    st.info(f"🛒 **TRONG GIỎ ĐANG CÓ: {len(st.session_state.photo_cart)} ẢNH** ĐÃ SẴN SÀNG.")
    if st.button("🗑️ Làm trống Giỏ hàng để chọn lại từ đầu"):
        st.session_state.photo_cart = {}
        st.rerun()

st.header("Bước 3: Tùy Chỉnh Layout")
mode = st.radio("Chọn chế độ:", ("Layout 1: Cơ bản (1-2 ảnh/trang)", "Layout 2: Tràn viền (Lưới tự động chia đều, nhiều ảnh)"))
align_mode = "2"
if "Layout 1" in mode:
    align_mode = st.radio("Căn lề cho ảnh (Layout 1):", ("1 - Trái", "2 - Giữa", "3 - Phải"))
    align_mode = align_mode[0]

st.header("Bước 4: Vị trí chèn & Thiết kế")
vitri_input = st.text_input("Chèn vào sau Slide số mấy? (Gõ 0 để chèn cuối cùng):", "0")

st.markdown("---")
use_blank = False
main_title, sub_title, cover_color = "", "", "#FFFFFF"
content_title, content_color = "", "#003366"
end_title, thankyou_color = "THANK YOU!", "#FFFFFF"
bg_cover_file, bg_content_file, bg_end_file = None, None, None

if not st.session_state.template_bytes:
    st.warning("⚠️ Bạn chưa tải file PowerPoint mẫu (ở Bước 1).")
    use_blank = st.checkbox("👉 Tích vào đây để BỎ QUA và TẠO FILE MỚI TINH")
    
    if use_blank:
        st.info("🎨 THIẾT KẾ CÁC TRANG CỦA FILE MỚI:")
        
        with st.expander("1️⃣ TRANG BÌA MỞ ĐẦU (Cover)", expanded=True):
            bg_cover_file = st.file_uploader("🖼️ Tải Ảnh Nền cho Bìa Mở Đầu:", type=['jpg', 'jpeg', 'png'], key="cover")
            main_title = st.text_input("Tiêu đề chính:", placeholder="VD: BÁO CÁO NGHIỆM THU")
            sub_title = st.text_input("Tiêu đề phụ:", placeholder="VD: Ngày 16/08/2026")
            col_grid1, col_color1 = st.columns([2, 1])
            with col_grid1: render_position_grid("cover_pos", "cv")
            with col_color1: cover_color = st.color_picker("🎨 Màu chữ bìa:", "#FFFFFF", key="c_col")

        with st.expander("2️⃣ TRANG NỘI DUNG (Các trang chứa ảnh)"):
            bg_content_file = st.file_uploader("🖼️ Tải Ảnh Nền chung cho các trang giữa:", type=['jpg', 'jpeg', 'png'], key="content_bg")
            content_title = st.text_input("Ghi chú góc trên trái (Sẽ tự động gạch chân):", placeholder="VD: HÌNH ẢNH THI CÔNG THỰC TẾ")
            content_color = st.color_picker("🎨 Màu chữ ghi chú:", "#003366", key="con_col")

        with st.expander("3️⃣ TRANG KẾT THÚC (Thank You)"):
            bg_end_file = st.file_uploader("🖼️ Tải Ảnh Nền cho Bìa Kết Thúc:", type=['jpg', 'jpeg', 'png'], key="end")
            end_title = st.text_input("Ghi chú kết thúc (Xóa trắng ô này nếu chỉ muốn hiện nền):", value="THANK YOU!")
            col_grid2, col_color2 = st.columns([2, 1])
            with col_grid2: render_position_grid("ty_pos", "ty")
            with col_color2: thankyou_color = st.color_picker("🎨 Màu chữ Kết thúc:", "#FFFFFF", key="t_col")

st.markdown("---")
# ==========================================
# CỤM NÚT TẠO FILE
# ==========================================
col_btn1, col_btn2 = st.columns(2)
with col_btn1:
    btn_pptx = st.button("🚀 TẠO FILE POWERPOINT", use_container_width=True, type="primary")
with col_btn2:
    btn_pdf = st.button("📄 TẠO FILE PDF NHANH", use_container_width=True)

if btn_pptx or btn_pdf:
    is_pdf = btn_pdf
    
    st.session_state.show_download_pptx = False
    st.session_state.show_download_pdf = False

    if is_pdf and st.session_state.template_bytes and not use_blank:
        st.error("⚠️ Xuất PDF trực tiếp chỉ hỗ trợ khi bác chọn 'TẠO FILE MỚI TINH'. Vui lòng tích vào ô Bỏ qua file mẫu ở Bước 4, hoặc chọn Xuất PowerPoint!")
    elif not st.session_state.template_bytes and not use_blank:
        st.error("⚠️ Vui lòng tải file mẫu ở Bước 1, HOẶC tích vào ô tạo file mới nhé!")
    elif not st.session_state.photo_cart:
        st.error("⚠️ Giỏ ảnh đang trống trơn! Bác chọn thêm ảnh ở Bước 2 nhé!")
    else:
        with st.spinner("Đang tự động dàn trang... Vui lòng đợi nhé..."):
            try:
                bg_cover_stream = process_bg_image(bg_cover_file) if use_blank else None
                bg_content_stream = process_bg_image(bg_content_file) if use_blank else None
                bg_end_stream = process_bg_image(bg_end_file) if use_blank else None
                
                r_cov, g_cov, b_cov = hex_to_rgb(cover_color)
                r_con, g_con, b_con = hex_to_rgb(content_color)
                r_ty, g_ty, b_ty = hex_to_rgb(thankyou_color)

                slide_w = Inches(13.3333333)
                slide_h = Inches(7.5)

                if not is_pdf:
                    if st.session_state.template_bytes:
                        prs = Presentation(io.BytesIO(st.session_state.template_bytes))
                    else:
                        prs = Presentation()
                        prs.slide_width = slide_w; prs.slide_height = slide_h
                    
                    tong_slide = len(prs.slides)
                    vi_tri_hien_tai = tong_slide 
                    if tong_slide > 0:
                        match = re.search(r'\d+', vitri_input)
                        if match:
                            trang_chon = int(match.group())
                            if 0 < trang_chon <= tong_slide: vi_tri_hien_tai = trang_chon
                    else: vi_tri_hien_tai = 0

                    try: slide_layout = prs.slide_layouts[6]
                    except: slide_layout = prs.slide_layouts[0] 

                pdf_slides = [] 

                # ==========================
                # TẠO TRANG BÌA MỞ ĐẦU
                # ==========================
                if use_blank and (main_title or sub_title or bg_cover_stream):
                    y_main = get_cover_y(st.session_state.cover_pos, slide_h, is_sub=False)
                    y_sub = get_cover_y(st.session_state.cover_pos, slide_h, is_sub=True)
                    align_cover = get_alignment(st.session_state.cover_pos)

                    if is_pdf:
                        slide_pdf = add_pdf_slide(pdf_slides, bg_cover_stream)
                        if main_title: draw_text_pdf(slide_pdf, main_title.upper(), Inches(0.5), y_main, slide_w - Inches(1), Inches(1), st.session_state.cover_pos, 44, cover_color)
                        if sub_title: draw_text_pdf(slide_pdf, sub_title, Inches(0.5), y_sub, slide_w - Inches(1), Inches(0.8), st.session_state.cover_pos, 24, cover_color)
                    else:
                        title_slide = prs.slides.add_slide(slide_layout)
                        if bg_cover_stream:
                            bg_cover_stream.seek(0)
                            title_slide.shapes.add_picture(bg_cover_stream, 0, 0, width=slide_w, height=slide_h)
                        
                        if main_title:
                            txBox = title_slide.shapes.add_textbox(Inches(0.5), y_main, slide_w - Inches(1), Inches(1))
                            tf = txBox.text_frame; tf.word_wrap = True
                            p = tf.paragraphs[0]; p.text = main_title.upper()
                            p.alignment = align_cover; p.font.size = Pt(44); p.font.bold = True
                            p.font.color.rgb = RGBColor(r_cov, g_cov, b_cov)

                        if sub_title:
                            txBox2 = title_slide.shapes.add_textbox(Inches(0.5), y_sub, slide_w - Inches(1), Inches(0.8))
                            tf2 = txBox2.text_frame; tf2.word_wrap = True
                            p2 = tf2.paragraphs[0]; p2.text = sub_title
                            p2.alignment = align_cover; p2.font.size = Pt(24)
                            p2.font.color.rgb = RGBColor(r_cov, g_cov, b_cov)
                        vi_tri_hien_tai = len(prs.slides)

                CACH_LE_TREN = Inches(0.9)  
                CACH_LE_DUOI = Inches(0.8)  
                CACH_LE_TRAI = Inches(0.2) 
                CACH_LE_PHAI = Inches(0.2) 
                GAP = Inches(0.12)          
                usable_w = slide_w - CACH_LE_TRAI - CACH_LE_PHAI
                usable_h = slide_h - CACH_LE_TREN - CACH_LE_DUOI

                # Lọc và sắp xếp ảnh
                image_data = []
                for key, img_info in st.session_state.photo_cart.items():
                    img_stream = io.BytesIO(img_info["bytes"])
                    try:
                        with Image.open(img_stream) as im:
                            im = ImageOps.exif_transpose(im)
                            width, height = im.size
                            is_portrait = height >= width
                            dt_str = "9999"
                            exif = im.getexif()
                            if exif: dt_str = exif.get(36867) or exif.get(306) or "9999"
                        img_stream.seek(0)
                        image_data.append({
                            'stream': img_stream, 'is_portrait': is_portrait, 
                            'w': width, 'h': height, 'name': img_info["name"], 'timestamp': str(dt_str)
                        })
                    except Exception: pass
                image_data.sort(key=lambda x: (x['timestamp'], x['name']))

                def add_content_title_pptx(slide_obj):
                    if use_blank and content_title:
                        tx = slide_obj.shapes.add_textbox(Inches(0.2), Inches(0.15), slide_w - Inches(0.4), Inches(0.6))
                        tf_c = tx.text_frame; p_c = tf_c.paragraphs[0]; p_c.text = content_title.upper()
                        p_c.font.size = Pt(22); p_c.font.bold = True; p_c.font.underline = True 
                        p_c.font.color.rgb = RGBColor(r_con, g_con, b_con)

                # ==========================
                # DÀN TRANG NỘI DUNG CHÍNH
                # ==========================
                if "Layout 1" in mode:
                    i = 0
                    while i < len(image_data):
                        current_img = image_data[i]
                        
                        if is_pdf:
                            slide_pdf = add_pdf_slide(pdf_slides, bg_content_stream)
                            if use_blank and content_title:
                                draw_text_pdf(slide_pdf, content_title.upper(), Inches(0.2), Inches(0.15), slide_w - Inches(0.4), Inches(0.6), "Trái", 22, content_color, underline=True)
                        else:
                            slide = prs.slides.add_slide(slide_layout) 
                            if use_blank and bg_content_stream:
                                bg_content_stream.seek(0)
                                slide.shapes.add_picture(bg_content_stream, 0, 0, width=slide_w, height=slide_h)
                            add_content_title_pptx(slide) 
                            move_slide(prs, len(prs.slides) - 1, vi_tri_hien_tai)
                        
                        if current_img['is_portrait'] and (i + 1 < len(image_data)) and image_data[i+1]['is_portrait']:
                            next_img = image_data[i+1]
                            r1 = current_img['w'] / current_img['h']
                            r2 = next_img['w'] / next_img['h']
                            test_w = usable_h * r1 + usable_h * r2
                            if test_w <= usable_w - GAP: final_h = usable_h
                            else: final_h = (usable_w - GAP) / (r1 + r2)
                                
                            final_w1 = final_h * r1; final_w2 = final_h * r2
                            block_w = final_w1 + GAP + final_w2
                            
                            if align_mode == '1': start_x = CACH_LE_TRAI
                            elif align_mode == '3': start_x = slide_w - CACH_LE_PHAI - block_w
                            else: start_x = CACH_LE_TRAI + (usable_w - block_w) / 2
                            start_y = CACH_LE_TREN + (usable_h - final_h) / 2
                            
                            if is_pdf:
                                add_image_pdf(slide_pdf, current_img['stream'], start_x, start_y, final_w1, final_h)
                                add_image_pdf(slide_pdf, next_img['stream'], start_x + final_w1 + GAP, start_y, final_w2, final_h)
                            else:
                                add_image_exact(slide, current_img['stream'], start_x, start_y, final_w1, final_h)
                                add_image_exact(slide, next_img['stream'], start_x + final_w1 + GAP, start_y, final_w2, final_h)
                            i += 2 
                        else:
                            r_img = current_img['w'] / current_img['h']
                            if usable_h * r_img <= usable_w:
                                f_h = usable_h; f_w = usable_h * r_img
                            else:
                                f_w = usable_w; f_h = usable_w / r_img
                                
                            if align_mode == '1': s_x = CACH_LE_TRAI
                            elif align_mode == '3': s_x = slide_w - CACH_LE_PHAI - f_w
                            else: s_x = CACH_LE_TRAI + (usable_w - f_w) / 2
                            s_y = CACH_LE_TREN + (usable_h - f_h) / 2
                            
                            if is_pdf: add_image_pdf(slide_pdf, current_img['stream'], s_x, s_y, f_w, f_h)
                            else: add_image_exact(slide, current_img['stream'], s_x, s_y, f_w, f_h)
                            i += 1
                        
                        if not is_pdf: vi_tri_hien_tai += 1

                elif "Layout 2" in mode:
                    landscapes = [img for img in image_data if not img['is_portrait']]
                    portraits = [img for img in image_data if img['is_portrait']]
                    smart_chunks = []
                    for c in partition_images(landscapes, 6): smart_chunks.append({'type': 'landscape', 'images': c})
                    for c in partition_images(portraits, 4): smart_chunks.append({'type': 'portrait', 'images': c})

                    for chunk_dict in smart_chunks:
                        chunk = chunk_dict['images']
                        n = len(chunk)
                        
                        if is_pdf:
                            slide_pdf = add_pdf_slide(pdf_slides, bg_content_stream)
                            if use_blank and content_title: draw_text_pdf(slide_pdf, content_title.upper(), Inches(0.2), Inches(0.15), slide_w - Inches(0.4), Inches(0.6), "Trái", 22, content_color, underline=True)
                        else:
                            slide = prs.slides.add_slide(slide_layout)
                            if use_blank and bg_content_stream:
                                bg_content_stream.seek(0)
                                slide.shapes.add_picture(bg_content_stream, 0, 0, width=slide_w, height=slide_h)
                            add_content_title_pptx(slide) 
                            move_slide(prs, len(prs.slides) - 1, vi_tri_hien_tai)

                        layout_rows = []
                        if chunk_dict['type'] == 'portrait': layout_rows = [chunk] 
                        else:
                            if n == 6: layout_rows = [chunk[0:3], chunk[3:6]]
                            elif n == 5: layout_rows = [chunk[0:3], chunk[3:5]]
                            elif n == 4: layout_rows = [chunk[0:2], chunk[2:4]]
                            else: layout_rows = [chunk]

                        if is_pdf: draw_adaptive_grid_pdf(slide_pdf, layout_rows, CACH_LE_TRAI, CACH_LE_TREN, usable_w, usable_h, GAP)
                        else: draw_adaptive_grid(slide, layout_rows, CACH_LE_TRAI, CACH_LE_TREN, usable_w, usable_h, GAP)
                        
                        if not is_pdf: vi_tri_hien_tai += 1

                # ==========================
                # TẠO TRANG KẾT THÚC
                # ==========================
                if use_blank and (end_title.strip() or bg_end_stream):
                    y_ty = get_ty_y(st.session_state.ty_pos, slide_h)
                    align_ty = get_alignment(st.session_state.ty_pos)

                    if is_pdf:
                        slide_pdf = add_pdf_slide(pdf_slides, bg_end_stream)
                        if end_title.strip(): draw_text_pdf(slide_pdf, end_title.upper(), Inches(0.5), y_ty, slide_w - Inches(1), Inches(1), st.session_state.ty_pos, 50, thankyou_color)
                    else:
                        ty_slide = prs.slides.add_slide(slide_layout)
                        if bg_end_stream:
                            bg_end_stream.seek(0)
                            ty_slide.shapes.add_picture(bg_end_stream, 0, 0, width=slide_w, height=slide_h)
                        move_slide(prs, len(prs.slides) - 1, vi_tri_hien_tai)

                        if end_title.strip():
                            txBox_ty = ty_slide.shapes.add_textbox(Inches(0.5), y_ty, slide_w - Inches(1), Inches(1))
                            tf_ty = txBox_ty.text_frame; tf_ty.word_wrap = True
                            p_ty = tf_ty.paragraphs[0]; p_ty.text = end_title.upper()
                            p_ty.alignment = align_ty; p_ty.font.size = Pt(50); p_ty.font.bold = True
                            p_ty.font.color.rgb = RGBColor(r_ty, g_ty, b_ty)

                output_stream = io.BytesIO()
                
                if is_pdf:
                    if len(pdf_slides) == 0:
                        st.error("⚠️ Không có gì để xuất PDF! Bác cần chèn ảnh hoặc bật chế độ Tạo file mới.")
                    else:
                        pdf_slides[0].save(output_stream, format="PDF", save_all=True, append_images=pdf_slides[1:])
                        st.session_state.final_pdf = output_stream.getvalue()
                        st.session_state.show_download_pdf = True
                else:
                    prs.save(output_stream)
                    st.session_state.final_pptx = output_stream.getvalue()
                    st.session_state.show_download_pptx = True

            except Exception as e:
                st.error(f"❌ Có lỗi xảy ra: {e}")

# ==========================================
# KHU VỰC HIỂN THỊ NÚT TẢI XUỐNG CỐ ĐỊNH
# ==========================================
if st.session_state.show_download_pptx and st.session_state.final_pptx:
    st.success("✅ Thành công mỹ mãn! File PowerPoint đã sẵn sàng.")
    st.download_button(
        label="📥 BẤM VÀO ĐÂY ĐỂ TẢI POWERPOINT VỀ MÁY",
        data=st.session_state.final_pptx,
        file_name="Report_Kiem_Tra.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        type="primary"
    )

if st.session_state.show_download_pdf and st.session_state.final_pdf:
    st.success("✅ Thành công mỹ mãn! File PDF đã sẵn sàng.")
    st.download_button(
        label="📥 BẤM VÀO ĐÂY ĐỂ TẢI PDF VỀ MÁY",
        data=st.session_state.final_pdf,
        file_name="Report_Kiem_Tra.pdf",
        mime="application/pdf",
        type="primary"
    )
